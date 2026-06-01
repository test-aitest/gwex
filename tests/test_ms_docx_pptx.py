"""docx / pptx parser の単体テスト（メモリ上で生成して検証）。"""

from __future__ import annotations

from io import BytesIO

import docx as pydocx
import pptx as pypptx
from pptx.util import Inches

from gwex.ir import Divider, Heading, Note, Table
from gwex.parsers import docx as docx_parser
from gwex.parsers import pptx as pptx_parser
from gwex.serializers import markdown as md


def test_docx_heading_paragraph_table():
    d = pydocx.Document()
    d.add_heading("見出し", level=1)
    p = d.add_paragraph("これは ")
    p.add_run("太字").bold = True
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "項目"
    t.cell(0, 1).text = "値"
    t.cell(1, 0).text = "A"
    t.cell(1, 1).text = "1"
    buf = BytesIO()
    d.save(buf)

    doc = docx_parser.parse(buf.getvalue())
    assert doc.source == "docx"
    out = md.render(doc)
    assert "# 見出し" in out
    assert "**太字**" in out
    assert "| 項目 | 値 |" in out
    assert "| A | 1 |" in out


def test_docx_list():
    d = pydocx.Document()
    d.add_paragraph("一", style="List Number")
    d.add_paragraph("二", style="List Number")
    buf = BytesIO()
    d.save(buf)
    out = md.render(docx_parser.parse(buf.getvalue()))
    assert "1. 一" in out
    assert "2. 二" in out


def test_pptx_slide_text_and_notes():
    prs = pypptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "タイトル"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "本文テキスト"
    slide.notes_slide.notes_text_frame.text = "発表者ノート"
    buf = BytesIO()
    prs.save(buf)

    doc = pptx_parser.parse(buf.getvalue())
    assert doc.source == "pptx"
    assert isinstance(doc.blocks[0], Divider)
    assert isinstance(doc.blocks[1], Heading)
    out = md.render(doc)
    assert "# Slide 1" in out
    assert "タイトル" in out
    assert "本文テキスト" in out
    assert "> **Note:**" in out
    assert "発表者ノート" in out
