"""Microsoft PowerPoint (.pptx) → IR。

pure 層: 与えられたバイト列を BytesIO で解析するだけ。各スライドを
`Divider` + `Heading(1)` "Slide N" に区切り、図形テキスト・表・ノートを写像する。
図形は (top, left) でソートして読み順を安定させる。
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from gwmd.ir import (
    Block,
    Divider,
    Document,
    Emph,
    Heading,
    Inline,
    Note,
    Paragraph,
    Strong,
    Table,
    Text,
)


def parse(data: bytes) -> Document:
    prs = Presentation(BytesIO(data))
    blocks: list[Block] = []
    for idx, slide in enumerate(prs.slides, start=1):
        blocks.append(Divider())
        blocks.append(Heading(level=1, inlines=[Text(value=f"Slide {idx}")]))
        for shape in sorted(slide.shapes, key=_topleft):
            blocks.extend(_shape(shape))
        notes = _notes(slide)
        if notes:
            blocks.append(Note(blocks=notes))
    return Document(title=None, source="pptx", blocks=blocks)


def _topleft(shape) -> tuple[int, int]:
    return (shape.top or 0, shape.left or 0)


def _shape(shape) -> list[Block]:
    if shape.has_table:
        return [_table(shape.table)]
    if shape.has_text_frame:
        return _text_frame_paragraphs(shape.text_frame)
    return []


def _text_frame_paragraphs(tf) -> list[Block]:
    out: list[Block] = []
    for para in tf.paragraphs:
        inlines = _inlines(para)
        if inlines:
            out.append(Paragraph(inlines=inlines))
    return out


def _inlines(para) -> list[Inline]:
    out: list[Inline] = []
    for run in para.runs:
        if run.text == "":
            continue
        node: Inline = Text(value=run.text)
        if run.font.italic:
            node = Emph(inlines=[node])
        if run.font.bold:
            node = Strong(inlines=[node])
        out.append(node)
    return out


def _table(table) -> Table:
    rows: list[list[list[Block]]] = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(_text_frame_paragraphs(cell.text_frame) or [Paragraph(inlines=[])])
        rows.append(cells)
    if not rows:
        return Table(header=None, rows=[])
    return Table(header=rows[0], rows=rows[1:])


def _notes(slide) -> list[Block]:
    if not slide.has_notes_slide:
        return []
    tf = slide.notes_slide.notes_text_frame
    return _text_frame_paragraphs(tf) if tf is not None else []
